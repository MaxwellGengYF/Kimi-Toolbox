"""Document-to-text extraction for ``ReadFile``.

Supported formats and the dedicated third-party parsers that back them:

=================  =========================  ========================================
Extension          Parser                     Notes
=================  =========================  ========================================
``.ipynb``         nbformat                   Official Jupyter notebook parser
``.docx``          python-docx                Paragraphs and tables, in document order
``.xlsx``/``.xlsm`` openpyxl                  Macro-enabled workbooks read as well
``.xls``           xlrd                       Legacy binary Excel format
``.pptx``          python-pptx                Slides, text frames and tables
``.pdf``           PyMuPDF (``pymupdf``)      Fast native PDF text extraction
=================  =========================  ========================================

Legacy binary ``.doc`` and ``.ppt`` files are intentionally unsupported:
python-docx / python-pptx cannot read them, and the usual fallback
(textract + antiword) is unmaintained and needs external system binaries.

Every parser is imported lazily inside its extractor so CLI startup cost is
unaffected. Malformed documents raise :class:`ExtractionError`; callers can
then fall back to normal text/binary handling.
"""

from __future__ import annotations

from pathlib import Path

__all__ = [
    "EXTRACTABLE_EXTENSIONS",
    "ExtractionError",
    "extract_document_text",
    "is_extractable_document",
]

EXTRACTABLE_EXTENSIONS = frozenset({".ipynb", ".docx", ".xlsx", ".xlsm", ".xls", ".pptx", ".pdf"})

# Row/column caps for spreadsheet extraction (LLM context is the bottleneck).
_MAX_ROWS_PER_SHEET = 5000
_MAX_COLS = 256
# PDFs can be huge; only the first pages are extracted.
_MAX_PDF_PAGES = 500


class ExtractionError(Exception):
    """Raised when a supported-looking document cannot be rendered as text."""


def _extension(path: str) -> str:
    ext = Path(path).suffix.lower()
    return ext if ext in EXTRACTABLE_EXTENSIONS else ""


def is_extractable_document(path: str) -> bool:
    return bool(_extension(path))


def extract_document_text(path: str) -> str:
    ext = _extension(path)
    if ext == ".ipynb":
        return _extract_notebook(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext in {".xlsx", ".xlsm"}:
        return _extract_xlsx(path)
    if ext == ".xls":
        return _extract_xls(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".pdf":
        return _extract_pdf(path)
    raise ExtractionError(f"Unsupported document type: {path!r}")


# ── .ipynb ─────────────────────────────────────────────────────────────────────


def _source_text(source) -> str:
    """Join a notebook cell source, which may be a str or a list of str."""
    if isinstance(source, str):
        return source
    if isinstance(source, list):
        return "".join(item for item in source if isinstance(item, str))
    return ""


def _extract_notebook(path: str) -> str:
    try:
        import nbformat
    except ImportError as exc:
        raise ExtractionError(f"nbformat is required for .ipynb extraction: {exc}") from exc
    try:
        nb = nbformat.read(path, as_version=4)
    except Exception as exc:
        raise ExtractionError(f"Not a valid notebook: {exc}") from exc
    cells = nb.get("cells") if isinstance(nb, dict) else None
    if not cells:
        raise ExtractionError("Notebook contains no cells")

    counts = {"markdown": 0, "code": 0, "raw": 0}
    labels = {"markdown": "Markdown", "code": "Code", "raw": "Raw"}
    out: list[str] = []
    for cell in cells:
        if not isinstance(cell, dict):
            continue
        typ = cell.get("cell_type")
        if typ not in labels:
            continue
        counts[typ] += 1
        suffix = f" {counts[typ]}" if typ != "raw" else ""
        out.extend(
            (
                f"# ── {labels[typ]} cell{suffix} ──",
                _source_text(cell.get("source", "")).rstrip("\n"),
                "",
            )
        )
    if not out:
        raise ExtractionError("Notebook contains no readable cells")
    return "\n".join(out).rstrip("\n") + "\n"


# ── .docx ──────────────────────────────────────────────────────────────────────


def _extract_docx(path: str) -> str:
    try:
        from docx import Document
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        document = Document(path)
        blocks: list[list[str]] = []
        for item in document.iter_inner_content():
            if isinstance(item, Paragraph):
                text = item.text
                blocks.append(text.split("\n") if "\n" in text else [text])
            elif isinstance(item, Table):
                blocks.append(["\t".join(cell.text for cell in row.cells) for row in item.rows])
        lines: list[str] = []
        for index, block in enumerate(blocks):
            if index:
                lines.append("")
            lines.extend(block)
        if not any(line.strip() for line in lines):
            raise ExtractionError("DOCX contains no extractable text")
        return "\n".join(lines).rstrip("\n") + "\n"
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(f"Not a valid DOCX: {exc}") from exc


# ── .xlsx / .xlsm ──────────────────────────────────────────────────────────────


def _extract_xlsx(path: str) -> str:
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=True, keep_links=False)
    except Exception as exc:
        raise ExtractionError(f"Not a valid XLSX: {exc}") from exc
    try:
        out: list[str] = []
        has_visible = False
        for sheet in workbook.worksheets:
            if sheet.sheet_state != "visible":
                continue
            has_visible = True
            rows = _iter_sheet_rows(sheet)
            out.append(f"# ── Sheet: {sheet.title} ──")
            out.extend("\t".join(row) for row in rows)
            if not rows:
                out.append("(empty)")
            out.append("")
        if not has_visible:
            raise ExtractionError("XLSX has no visible sheets")
        return "\n".join(out).rstrip("\n") + "\n"
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(f"Failed to read XLSX: {exc}") from exc
    finally:
        workbook.close()


def _iter_sheet_rows(sheet) -> list[list[str]]:
    rows: list[list[str]] = []
    for row in sheet.iter_rows(values_only=True):
        if len(rows) >= _MAX_ROWS_PER_SHEET:
            break
        cells = [_cell_text(value) for value in row[:_MAX_COLS]]
        while cells and not cells[-1]:
            cells.pop()
        rows.append(cells)
    while rows and not any(rows[-1]):
        rows.pop()
    return rows


# ── .xls ───────────────────────────────────────────────────────────────────────


def _extract_xls(path: str) -> str:
    try:
        import xlrd

        # Default (on_demand=False) reads the whole book up front and closes
        # the file handle before returning.
        book = xlrd.open_workbook(path)
    except Exception as exc:
        raise ExtractionError(f"Not a valid XLS: {exc}") from exc
    out: list[str] = []
    has_visible = False
    for sheet in book.sheets():
        if sheet.visibility != 0:
            continue
        has_visible = True
        rows = _iter_xls_rows(sheet)
        out.append(f"# ── Sheet: {sheet.name} ──")
        out.extend("\t".join(row) for row in rows)
        if not rows:
            out.append("(empty)")
        out.append("")
    if not has_visible:
        raise ExtractionError("XLS has no visible sheets")
    return "\n".join(out).rstrip("\n") + "\n"


def _iter_xls_rows(sheet) -> list[list[str]]:
    rows: list[list[str]] = []
    limit_r = min(sheet.nrows, _MAX_ROWS_PER_SHEET)
    limit_c = min(sheet.ncols, _MAX_COLS)
    for r in range(limit_r):
        cells = [_cell_text(sheet.cell_value(r, c)) for c in range(limit_c)]
        while cells and not cells[-1]:
            cells.pop()
        rows.append(cells)
    while rows and not any(rows[-1]):
        rows.pop()
    return rows


def _cell_text(value) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


# ── .pptx ──────────────────────────────────────────────────────────────────────


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(path)
        out: list[str] = []
        has_text = False
        for index, slide in enumerate(prs.slides, 1):
            lines = _slide_text_lines(slide.shapes)
            if lines:
                has_text = True
            out.append(f"# ── Slide {index} ──")
            out.extend(lines if lines else ["(empty)"])
            out.append("")
        if not has_text:
            raise ExtractionError("PPTX contains no extractable text")
        return "\n".join(out).rstrip("\n") + "\n"
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(f"Not a valid PPTX: {exc}") from exc


def _slide_text_lines(shapes) -> list[str]:
    lines: list[str] = []
    for shape in shapes:
        group = getattr(shape, "shapes", None)
        if group is not None:
            lines.extend(_slide_text_lines(group))
            continue
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is not None:
            lines.extend(paragraph.text for paragraph in text_frame.paragraphs)
        if shape.has_table:
            table = getattr(shape, "table", None)
            if table is not None:
                lines.extend("\t".join(cell.text for cell in row.cells) for row in table.rows)
    return lines


# ── .pdf ───────────────────────────────────────────────────────────────────────


def _extract_pdf(path: str) -> str:
    try:
        import pymupdf
    except ImportError as exc:
        raise ExtractionError(f"PyMuPDF is required for .pdf extraction: {exc}") from exc
    try:
        document = pymupdf.open(path)
    except Exception as exc:
        raise ExtractionError(f"Not a valid PDF: {exc}") from exc
    try:
        out: list[str] = []
        has_text = False
        for index, page in enumerate(document, 1):
            if index > _MAX_PDF_PAGES:
                break
            text = page.get_text("text").rstrip("\n")
            if text:
                has_text = True
            out.append(f"# ── Page {index} ──")
            out.append(text if text else "(empty)")
            out.append("")
        if not has_text:
            raise ExtractionError("PDF contains no extractable text")
        return "\n".join(out).rstrip("\n") + "\n"
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(f"Failed to extract PDF text: {exc}") from exc
    finally:
        document.close()

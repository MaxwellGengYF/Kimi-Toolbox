"""Tests for document extraction in ReadFile (.ipynb/.docx/.xlsx/.xlsm/.xls/.pptx/.pdf).

Fixtures are built with the same third-party parsers that read_extract uses
(python-docx, openpyxl, python-pptx, xlwt, PyMuPDF) so the tests exercise the
real file formats end to end.
"""

from __future__ import annotations

import zipfile
from pathlib import Path

import orjson
import pytest
from kaos.path import KaosPath

from kimi_cli.tools.file import read as read_module
from kimi_cli.tools.file.read import Params, ReadFile
from kimi_cli.tools.file.read_extract import (
    EXTRACTABLE_EXTENSIONS,
    ExtractionError,
    extract_document_text,
    is_extractable_document,
)

NB_JSON = {
    "nbformat": 4,
    "nbformat_minor": 5,
    "metadata": {},
    "cells": [
        {
            "id": "cell-1",
            "cell_type": "markdown",
            "metadata": {},
            "source": ["# Title\n", "Some *text*."],
        },
        {"id": "cell-2", "cell_type": "code", "metadata": {}, "source": "print('hello')"},
        {"id": "cell-3", "cell_type": "raw", "metadata": {}, "source": "raw content"},
    ],
}

# v3 notebooks use `worksheets` + code cells with an `input` field; nbformat
# upgrades them to v4 internally.
NB_JSON_V3 = {
    "nbformat": 3,
    "nbformat_minor": 0,
    "metadata": {},
    "worksheets": [
        {
            "metadata": {},
            "cells": [
                {
                    "cell_type": "code",
                    "metadata": {},
                    "input": ["x = 1"],
                    "outputs": [],
                },
            ],
        }
    ],
}


def _write_zip(path: KaosPath, files: dict[str, str | bytes]) -> None:
    """Write a minimal zip archive from {entry_name: content}."""
    with zipfile.ZipFile(str(path), "w") as zf:
        for name, content in files.items():
            data = content.encode("utf-8") if isinstance(content, str) else content
            zf.writestr(name, data)


def _write_bytes(path: KaosPath, data: bytes) -> None:
    """Synchronous write for module-level (non-async) tests."""
    Path(str(path)).write_bytes(data)


# ── fixture builders (same libs as the extractor) ─────────────────────────────


def _build_docx(path: KaosPath, *, with_table: bool = False) -> None:
    from docx import Document

    document = Document()
    paragraph = document.add_paragraph()
    paragraph.add_run("Hello\tWorld")
    document.add_paragraph("Second para")
    if with_table:
        table = document.add_table(rows=1, cols=2)
        table.rows[0].cells[0].text = "H1"
        table.rows[0].cells[1].text = "H2"
    document.save(str(path))


def _build_xlsx(path: KaosPath, *, hidden: bool = False) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["Name", 42])
    sheet.append(["Alice"])
    if hidden:
        hidden_sheet = workbook.create_sheet("HiddenSheet")
        hidden_sheet.sheet_state = "hidden"
        hidden_sheet.append(["secret"])
    workbook.save(str(path))


def _build_pptx(path: KaosPath) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text = "Title line"
    bullet = box.text_frame.add_paragraph()
    bullet.text = "Bullet line"

    slide2 = prs.slides.add_slide(blank)
    graphic = slide2.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(4), Inches(1))
    graphic.table.cell(0, 0).text = "A"
    graphic.table.cell(0, 1).text = "B"
    prs.save(str(path))


def _build_xls(path: KaosPath) -> None:
    import xlwt

    book = xlwt.Workbook()
    sheet = book.add_sheet("Old")
    sheet.write(0, 0, "Name")
    sheet.write(0, 1, 42)
    sheet.write(1, 0, "Bob")
    book.save(str(path))


def _build_pdf(path: KaosPath, pages: int = 2) -> None:
    import pymupdf

    doc = pymupdf.open()
    for index in range(pages):
        page = doc.new_page()
        page.insert_text((72, 72), f"PDF text page {index + 1}")
    doc.save(str(path))
    doc.close()


# ── extension routing ─────────────────────────────────────────────────────────


def test_is_extractable_document():
    assert set(EXTRACTABLE_EXTENSIONS) == {
        ".ipynb",
        ".docx",
        ".xlsx",
        ".xlsm",
        ".xls",
        ".pptx",
        ".pdf",
    }
    assert is_extractable_document("notes.ipynb")
    assert is_extractable_document("doc.DOCX")
    assert is_extractable_document("a/b/c.xlsx")
    assert is_extractable_document("macro.xlsm")
    assert is_extractable_document("legacy.xls")
    assert is_extractable_document("deck.PPTX")
    assert is_extractable_document("paper.pdf")
    assert not is_extractable_document("notes.txt")
    assert not is_extractable_document("noext")
    assert not is_extractable_document("archive.zip")
    assert not is_extractable_document("old.doc")
    assert not is_extractable_document("old.ppt")


# ── .ipynb ────────────────────────────────────────────────────────────────────


def test_extract_notebook_markdown_and_code_cells(temp_work_dir: KaosPath):
    nb = temp_work_dir / "nb.ipynb"
    _write_bytes(nb, orjson.dumps(NB_JSON))

    text = extract_document_text(str(nb))
    assert text == (
        "# ── Markdown cell 1 ──\n"
        "# Title\n"
        "Some *text*.\n"
        "\n"
        "# ── Code cell 1 ──\n"
        "print('hello')\n"
        "\n"
        "# ── Raw cell ──\n"
        "raw content\n"
    )


def test_extract_notebook_legacy_v3_worksheets(temp_work_dir: KaosPath):
    """Legacy v3 notebooks (worksheets + code `input` cells) are upgraded."""
    nb = temp_work_dir / "legacy.ipynb"
    _write_bytes(nb, orjson.dumps(NB_JSON_V3))
    text = extract_document_text(str(nb))
    assert text == "# ── Code cell 1 ──\nx = 1\n"


def test_extract_notebook_no_cells_raises(temp_work_dir: KaosPath):
    nb = temp_work_dir / "empty.ipynb"
    _write_bytes(
        nb, orjson.dumps({"nbformat": 4, "nbformat_minor": 5, "metadata": {}, "cells": []})
    )
    with pytest.raises(ExtractionError, match="no cells"):
        extract_document_text(str(nb))


def test_extract_notebook_bad_json_raises(temp_work_dir: KaosPath):
    nb = temp_work_dir / "broken.ipynb"
    _write_bytes(nb, b"{this is not json")
    with pytest.raises(ExtractionError, match="Not a valid notebook"):
        extract_document_text(str(nb))


# ── .docx ─────────────────────────────────────────────────────────────────────


def test_extract_docx_paragraphs_tabs_breaks(temp_work_dir: KaosPath):
    docx = temp_work_dir / "doc.docx"
    _build_docx(docx)

    text = extract_document_text(str(docx))
    assert text == "Hello\tWorld\n\nSecond para\n"


def test_extract_docx_table(temp_work_dir: KaosPath):
    docx = temp_work_dir / "table.docx"
    _build_docx(docx, with_table=True)

    text = extract_document_text(str(docx))
    assert text == "Hello\tWorld\n\nSecond para\n\nH1\tH2\n"


def test_extract_bad_docx_raises(temp_work_dir: KaosPath):
    bad = temp_work_dir / "bad.docx"
    _write_bytes(bad, b"this is not a zip file at all")
    with pytest.raises(ExtractionError, match="Not a valid DOCX"):
        extract_document_text(str(bad))


# ── .xlsx / .xlsm ─────────────────────────────────────────────────────────────


def test_extract_xlsx_sheets_and_values(temp_work_dir: KaosPath):
    xlsx = temp_work_dir / "book.xlsx"
    _build_xlsx(xlsx, hidden=True)

    text = extract_document_text(str(xlsx))
    assert text == "# ── Sheet: Data ──\nName\t42\nAlice\n"


def test_extract_xlsm_same_as_xlsx(temp_work_dir: KaosPath):
    xlsm = temp_work_dir / "book.xlsm"
    _build_xlsx(xlsm)

    text = extract_document_text(str(xlsm))
    assert text == "# ── Sheet: Data ──\nName\t42\nAlice\n"


def test_extract_xlsx_empty_sheet_shows_placeholder(temp_work_dir: KaosPath):
    from openpyxl import Workbook

    xlsx = temp_work_dir / "empty.xlsx"
    workbook = Workbook()
    workbook.save(str(xlsx))

    text = extract_document_text(str(xlsx))
    assert text == "# ── Sheet: Sheet ──\n(empty)\n"


def test_extract_bad_xlsx_raises(temp_work_dir: KaosPath):
    bad = temp_work_dir / "bad.xlsx"
    _write_bytes(bad, b"this is not a zip file at all")
    with pytest.raises(ExtractionError, match="Not a valid XLSX"):
        extract_document_text(str(bad))


# ── .xls ──────────────────────────────────────────────────────────────────────


def test_extract_xls_legacy_sheets(temp_work_dir: KaosPath):
    xls = temp_work_dir / "legacy.xls"
    _build_xls(xls)

    text = extract_document_text(str(xls))
    assert text == "# ── Sheet: Old ──\nName\t42\nBob\n"


# ── .pptx ─────────────────────────────────────────────────────────────────────


def test_extract_pptx_slides_text_and_table(temp_work_dir: KaosPath):
    pptx = temp_work_dir / "deck.pptx"
    _build_pptx(pptx)

    text = extract_document_text(str(pptx))
    assert text == ("# ── Slide 1 ──\nTitle line\nBullet line\n\n# ── Slide 2 ──\nA\tB\n")


def test_extract_bad_pptx_raises(temp_work_dir: KaosPath):
    bad = temp_work_dir / "bad.pptx"
    _write_bytes(bad, b"this is not a zip file at all")
    with pytest.raises(ExtractionError, match="Not a valid PPTX"):
        extract_document_text(str(bad))


# ── .pdf ──────────────────────────────────────────────────────────────────────


def test_extract_pdf_pages(temp_work_dir: KaosPath):
    pdf = temp_work_dir / "paper.pdf"
    _build_pdf(pdf, pages=2)

    text = extract_document_text(str(pdf))
    assert text == ("# ── Page 1 ──\nPDF text page 1\n\n# ── Page 2 ──\nPDF text page 2\n")


def test_extract_bad_pdf_raises(temp_work_dir: KaosPath):
    bad = temp_work_dir / "bad.pdf"
    _write_bytes(bad, b"this is not a pdf at all")
    with pytest.raises(ExtractionError, match="Not a valid PDF"):
        extract_document_text(str(bad))


# ── tool-level tests ──────────────────────────────────────────────────────────


async def test_read_ipynb_via_tool(read_file_tool: ReadFile, temp_work_dir: KaosPath):
    nb = temp_work_dir / "nb.ipynb"
    await nb.write_bytes(orjson.dumps(NB_JSON))
    display_path = str(nb).replace("\\", "/")

    result = await read_file_tool(Params(path=str(nb)))
    assert not result.is_error
    assert "## Markdown cell 1" in result.output
    assert "## Code cell 1" in result.output
    assert "print('hello')" in result.output
    assert "raw content" in result.output
    assert result.message.startswith(
        "11 lines read from file starting from line 1. Total lines in file: 11. End of file reached."
    )
    assert "(extracted from .ipynb document)" in result.message
    assert result.message.endswith(f" Path: {display_path}")


async def test_read_docx_via_tool(read_file_tool: ReadFile, temp_work_dir: KaosPath):
    docx = temp_work_dir / "doc.docx"
    _build_docx(docx)
    display_path = str(docx).replace("\\", "/")

    result = await read_file_tool(Params(path=str(docx)))
    assert not result.is_error
    assert "Hello\tWorld" in result.output
    assert "Second para" in result.output
    assert result.message.startswith(
        "2 lines read from file starting from line 1. Total lines in file: 2. End of file reached."
    )
    assert "(extracted from .docx document)" in result.message
    assert result.message.endswith(f" Path: {display_path}")


async def test_read_xlsx_via_tool(read_file_tool: ReadFile, temp_work_dir: KaosPath):
    xlsx = temp_work_dir / "book.xlsx"
    _build_xlsx(xlsx, hidden=True)
    display_path = str(xlsx).replace("\\", "/")

    result = await read_file_tool(Params(path=str(xlsx)))
    assert not result.is_error
    assert "## Sheet: Data" in result.output
    assert "| Name | 42 |" in result.output
    assert "Alice" in result.output
    # Hidden sheet is skipped.
    assert "HiddenSheet" not in result.output
    assert "secret" not in result.output
    assert "(extracted from .xlsx document)" in result.message
    assert result.message.endswith(f" Path: {display_path}")


async def test_read_pptx_via_tool(read_file_tool: ReadFile, temp_work_dir: KaosPath):
    pptx = temp_work_dir / "deck.pptx"
    _build_pptx(pptx)

    result = await read_file_tool(Params(path=str(pptx)))
    assert not result.is_error
    assert "## Slide 1" in result.output
    assert "Title line" in result.output
    assert "| A | B |" in result.output
    assert "(extracted from .pptx document)" in result.message
    assert "Path:" in result.message


async def test_read_pdf_via_tool(read_file_tool: ReadFile, temp_work_dir: KaosPath):
    pdf = temp_work_dir / "paper.pdf"
    _build_pdf(pdf, pages=1)

    result = await read_file_tool(Params(path=str(pdf)))
    assert not result.is_error
    assert "## Page 1" in result.output
    assert "PDF text page 1" in result.output
    assert "(extracted from .pdf document)" in result.message
    assert "Path:" in result.message


async def test_read_malformed_docx_bad_zip(read_file_tool: ReadFile, temp_work_dir: KaosPath):
    bad = temp_work_dir / "bad.docx"
    await bad.write_bytes(b"\x00\x01\x02 this is not a zip")

    result = await read_file_tool(Params(path=str(bad)))
    assert result.is_error
    assert "could not be extracted" in result.message
    assert result.brief == "Document extraction failed"


async def test_read_docx_bad_package_zip(read_file_tool: ReadFile, temp_work_dir: KaosPath):
    """A zip that is not a valid DOCX package fails cleanly."""
    docx = temp_work_dir / "nodoc.docx"
    _write_zip(docx, {"word/styles.xml": "<styles/>"})

    result = await read_file_tool(Params(path=str(docx)))
    assert result.is_error
    assert "could not be extracted" in result.message
    assert "Not a valid DOCX" in result.message
    assert result.brief == "Document extraction failed"


async def test_read_ipynb_no_cells(read_file_tool: ReadFile, temp_work_dir: KaosPath):
    nb = temp_work_dir / "empty.ipynb"
    await nb.write_bytes(orjson.dumps({"cells": []}))

    result = await read_file_tool(Params(path=str(nb)))
    assert result.is_error
    assert "could not be extracted" in result.message
    assert result.brief == "Document extraction failed"


async def test_read_extract_size_guard(
    read_file_tool: ReadFile, temp_work_dir: KaosPath, monkeypatch
):
    """Documents over MAX_EXTRACT_BYTES are refused without extraction."""
    docx = temp_work_dir / "big.docx"
    _build_docx(docx)
    monkeypatch.setattr(read_module, "MAX_EXTRACT_BYTES", 4)

    result = await read_file_tool(Params(path=str(docx)))
    assert result.is_error
    assert "larger than" in result.message
    assert "cannot be extracted" in result.message
    assert result.brief == "File not readable"


async def test_read_extract_pagination(read_file_tool: ReadFile, temp_work_dir: KaosPath):
    """char_offset/max_char pagination applies to extracted text."""
    nb = temp_work_dir / "nb.ipynb"
    await nb.write_bytes(orjson.dumps(NB_JSON))

    full = await read_file_tool(Params(path=str(nb), char_offset=0, max_char=200000))
    p1 = await read_file_tool(Params(path=str(nb), char_offset=0, max_char=30))
    p2 = await read_file_tool(Params(path=str(nb), char_offset=30, max_char=30))

    assert not full.is_error and not p1.is_error and not p2.is_error
    assert len(p1.output) <= 30
    assert len(p2.output) <= 30
    assert p1.output == full.output[0:30]
    assert p2.output == full.output[30:60]
    assert p1.output != p2.output
